import os
import streamlit as st
import google.generativeai as genai
import google.ai.generativelanguage as glm
from PIL import Image
from dotenv import load_dotenv
import streamlit_authenticator as stauth
import pickle
from pathlib import Path
import docx2txt
import pandas as pd
import PyPDF2
import pptx
import openai  # Import the OpenAI library

# Load the .env file
load_dotenv()

# Access the API keys
google_api_key = os.getenv('GOOGLE_API_KEY')
openai_api_key = os.getenv('OPENAI_API_KEY')

# Set API keys in the respective libraries
openai.api_key = openai_api_key

# Ensure the Google API key is set for Generative AI if needed
genai.configure(api_key=google_api_key)

st.set_page_config(page_title="XrossGlobal Financial - File Analysis", page_icon="📸", layout="centered", initial_sidebar_state='collapsed')
col1, col2, col3 = st.columns([1,2,1])

# Function to translate text using OpenAI's new API
def translate_with_chatgpt(text, target_language):
    prompt = f"Translate the following text to {target_language}:\n\n{text}"
    response = openai.chat.completions.create(  # Corrected API call
        model="gpt-4-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that can communicate in many languages."},
            {"role": "user", "content": prompt}
        ],
        temperature=0
    )
    return response.choices[0].message.content.strip()

# Language selection menu
languages = [
    "English", "Afrikaans", "Albanian", "Arabic", "Armenian", "Basque", "Belarusian",
    "Bengali", "Bosnian", "Bulgarian", "Catalan", "Chinese (Simplified)",
    "Chinese (Traditional)", "Croatian", "Czech", "Danish", "Dutch", "Estonian",
    "Filipino/Tagalog", "Finnish", "French", "Galician", "Georgian", "German",
    "Greek", "Gujarati", "Hebrew", "Hindi", "Hungarian", "Icelandic", "Irish",
    "Italian", "Japanese", "Kannada", "Korean", "Latvian", "Lithuanian",
    "Macedonian", "Malay/Indonesian", "Malayalam", "Maltese", "Marathi",
    "Mongolian", "Nepali", "Norwegian", "Pashto", "Persian (Farsi)", "Polish",
    "Portuguese", "Punjabi", "Romanian", "Russian", "Serbian", "Sinhalese",
    "Slovak", "Slovenian", "Spanish", "Swahili", "Swedish", "Tamil", "Telugu",
    "Thai", "Turkish", "Ukrainian", "Urdu", "Vietnamese", "Welsh", "Yoruba", "Zulu"
]

chosen_language = st.selectbox("Please choose a language:", languages)

# Texts to be translated
welcome_text = "Welcome *{name}*"
st.markdown(
    """
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Digitalo+Tech&display=swap');
    .custom-font {
        font-family: 'Michroma', sans-serif;
        text-align: center;
        font-size: 35px;
        color: maroon;
    }
    </style>
    <div class="custom-font">
        <strong>XrossGlobal Financial</strong>
    </div>
    """,
    unsafe_allow_html=True
)

st.markdown(
    """
    <div style="text-align: center; font-size: 20px;">
        <strong>AI powered FinTech & RegTech solutions for Banks & other Corporate Clients</strong>
    </div>
    """,
    unsafe_allow_html=True
)

header_text = "File Analysis"

# URL of the image
image_url = "https://i0.wp.com/xrossglobal.com/wp-content/uploads/2021/08/Russ-Chidy-stochastics-Final.jpg?resize=323%2C330&ssl=1"

# Display the image using st.image
#st.image(image_url, caption='XrossGlobal', width=100)

# Center the image using HTML and CSS
st.markdown(
    f"""
    <div style="text-align: center;">
        <img src="{image_url}" alt="XrossGlobal - AI for global users" width="100">
        <p>XrossGlobal - AI for global users</p>
    </div>
    """,
    unsafe_allow_html=True
)

subheader_text = "Please upload your files now for Analysis"
upload_text = "Upload PDF, Word, Excel, PowerPoint, CSV, or Images for analysis."

# Translate texts if the chosen language is not English
if chosen_language != "English":
    welcome_text = translate_with_chatgpt(welcome_text, chosen_language)
    header_text = translate_with_chatgpt(header_text, chosen_language)
    subheader_text = translate_with_chatgpt(subheader_text, chosen_language)
    upload_text = translate_with_chatgpt(upload_text, chosen_language)

# User Authentication Setup
names = ["Russ Chidy", "Tini Chidy", "Pana Chidy","Lykaa Nguyen","Vincent Le"]
usernames = ["russchidy", "tinichidy", "panachidy","lykaa","vincent"]
passwords = ["Tinman24", "Tinashe1", "Panashe1","Lykaa2024","Vinny_sgp"]

hashed_passwords = stauth.Hasher(passwords).generate()


# Load hashed passwords
file_path = Path(__file__).parent / "hashed_pw.pkl"
with file_path.open("wb") as file:
  pickle.dump(hashed_passwords, file)
    #hashed_passwords = pickle.load(file)

authenticator = stauth.Authenticate(
    credentials={"usernames": {usernames[i]: {"name": names[i], "password": hashed_passwords[i]} for i in range(len(names))}},
    cookie_name="xrossglobal_cookie",
    cookie_key="Tinman24",
    cookie_expiry_days=30
)

name, authentication_status, username = authenticator.login(location='main')

if authentication_status == False:
    st.error('Username/password is incorrect')

elif authentication_status == None:
    st.warning('Please enter your username and password')

elif authentication_status:
      st.write(welcome_text.format(name=name))
      st.header(header_text)
      st.subheader(subheader_text)
      st.write(upload_text)

      uploaded_file = st.file_uploader(
          "Choose a file",
          type=['jpg', 'jpeg', 'png', 'pdf', 'docx', 'xlsx', 'pptx', 'csv', 'tiff', 'webp']
      )

      text = None
      bytes_data = None

      if uploaded_file:
          file_extension = uploaded_file.name.split(".")[-1].lower()

          if file_extension in ['jpg', 'jpeg', 'png', 'tiff', 'webp']:
              # Handle image files
              image = Image.open(uploaded_file)
              st.image(image, caption='Uploaded Image', use_column_width=True)
              bytes_data = uploaded_file.getvalue()

              # Debugging: Check the content of bytes_data
              st.write(f"Bytes data length: {len(bytes_data)}")

              # Analyze Button
              generate = st.button("Click to Analyze!")

              if generate:
                  try:
                      model = genai.GenerativeModel('gemini-1.5-flash')

                      # Generate content using the AI model
                      response = model.generate_content(glm.Content(parts=[
                          glm.Part(text="You are acting as an expert in processes and Risk Management for Financial institutions. If the uploaded image is a math problem, give step-by-step guidance on how to solve it. "
                                        "For other images, provide a detailed description. Briefly list summary bullet points for any BIS or IOSCO frameworks only if relevant."),
                          glm.Part(inline_data=glm.Blob(mime_type='image/jpeg', data=bytes_data)),
                      ]), stream=True)

                      response.resolve()

                      # Check for safety ratings
                      if hasattr(response, 'safety_ratings') and response.safety_ratings:
                          safety_concerns = [rating for rating in response.safety_ratings if rating != 'SAFE']
                          if safety_concerns:
                              st.error("The image could not be processed due to safety concerns.")
                          else:
                              response_text = response.text
                              if chosen_language != "English":
                                  response_text = translate_with_chatgpt(response_text, chosen_language)
                              st.write(response_text)
                      else:
                          response_text = response.text
                          if chosen_language != "English":
                              response_text = translate_with_chatgpt(response_text, chosen_language)
                          st.write(response_text)

                      # Store the response text for later use
                      st.session_state['image_description'] = response_text

                  except Exception as e:
                      st.error(f"An error occurred: {str(e)}")

          else:
              # Handle non-image files
              if file_extension == 'pdf':
                  # Handle PDF files
                  pdf_reader = PyPDF2.PdfReader(uploaded_file)
                  text = ""
                  for page in pdf_reader.pages:
                      text += page.extract_text()

              elif file_extension == 'docx':
                  # Handle Word files
                  text = docx2txt.process(uploaded_file)

              elif file_extension == 'xlsx':
                  # Handle Excel files
                  df = pd.read_excel(uploaded_file)
                  text = df.to_string()

              elif file_extension == 'csv':
                  # Handle CSV files
                  try:
                      df = pd.read_csv(uploaded_file, on_bad_lines='warn')
                      text = df.to_string()
                  except pd.errors.ParserError as e:
                      st.error(f"An error occurred while processing the CSV file: {str(e)}")

              elif file_extension == 'pptx':
                  # Handle PowerPoint files
                  ppt = pptx.Presentation(uploaded_file)
                  text = ""
                  for slide in ppt.slides:
                      for shape in slide.shapes:
                          if hasattr(shape, "text"):
                              text += shape.text + "\n"

              # Display the extracted text
              if text:
                  st.write("Extracted Text:")
                  st.write(text)

              # Translate the text into the chosen language if necessary
              if chosen_language != "English":
                  translated_text = translate_with_chatgpt(text, chosen_language)
                  st.write(f"You have chosen to communicate in :{chosen_language}", chosen_language)
                  st.write(f"Translated Text ({chosen_language}):")
                  st.write(translated_text)
              else:
                  translated_text = text

          # Ask user if they want to ask questions about the document or image
          if chosen_language != "English":
              subheader_text_translated = translate_with_chatgpt(f"Ask questions about the document or image in {chosen_language}:", chosen_language)
              text_input_label_translated = translate_with_chatgpt(f"Type your question in {chosen_language}:", chosen_language)
          else:
              subheader_text_translated = f"Ask questions about the document or image in {chosen_language}:"
              text_input_label_translated = f"Type your question in {chosen_language}:"

          # Display the translated subheader and text input
          st.subheader(subheader_text_translated)
          user_question = st.text_input(text_input_label_translated)

          if user_question:
              # Use the LLM to answer the question in the selected language
              if 'image_description' in st.session_state:
                  context_text = st.session_state['image_description']
              else:
                  context_text = translated_text

              prompt = f"Answer the question in as much detail as possible based on the documents or image uploaded. Provide your response in the {chosen_language}:\n\n{context_text}\n\nAnswer the following question. If the uploaded document or image contains a maths formula or maths problem give a detailed step-by-step explanation on how to derive the answer:\n{user_question}"

              try:
                  model = genai.GenerativeModel('gemini-1.5-flash')

                  # Generate content using the AI model
                  response = model.generate_content(glm.Content(parts=[
                      glm.Part(text=prompt),
                  ]), stream=True)

                  response.resolve()

                  # Check for safety ratings
                  if hasattr(response, 'safety_ratings') and response.safety_ratings:
                      safety_concerns = [rating for rating in response.safety_ratings if rating != 'SAFE']
                      if safety_concerns:
                          st.error("The content could not be processed due to safety concerns.")
                      else:
                          response_text = response.text
                          if chosen_language != "English":
                              response_text = translate_with_chatgpt(response_text, chosen_language)
                          st.write(response_text)
                  else:
                      response_text = response.text
                      if chosen_language != "English":
                          response_text = translate_with_chatgpt(response_text, chosen_language)
                      st.write(response_text)

              except Exception as e:
                  st.error(f"An error occurred: {str(e)}")

      # Sidebar and Logout
      if authentication_status:
          authenticator.logout('Logout', 'sidebar')
          st.sidebar.title(f'Welcome *{name}*')
          user_photo = st.sidebar.camera_input("Take user photo - this is useful for high security scenarios such as High-Value payments that require face recognition for user authentication.")
          if user_photo:
              st.sidebar.image(user_photo)
          st.sidebar.header('Menu to come here')
